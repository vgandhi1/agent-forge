"""Binary-document tool layer for the contract-revision POC.

Why this exists: AgentForge's stock ``write_file`` / ``edit_file`` are plain-text only.
``.docx`` / ``.pptx`` are zipped-XML bundles — writing them as text corrupts the file.
These tools open the bundle with python-docx / python-pptx, mutate the object model, and
save it back, so formatting, styles, headers, and tables survive the edit.

Mirrors the safety model of ``core.artifact_store``: every path is validated to stay under
a single root (no absolute paths, no ``..`` traversal) before any read or write. The model
never gets arbitrary filesystem access — only relative paths under the contract workspace.

The handlers return plain strings (same contract as AgentForge tool handlers:
``(tool_input: dict) -> str``) so they drop straight into ``run_tool_loop``-style dispatch.
"""

from __future__ import annotations

import re
from pathlib import Path

from docx import Document
from pptx import Presentation

# Same conservative relative-path shape AgentForge uses in core/artifact_store.py.
_SAFE_REL = re.compile(r"^[A-Za-z0-9._\-/ ]+$")


class DocStore:
    """Sandbox a directory and expose binary-safe docx/pptx read + edit operations.

    All public methods take a *relative* path and resolve it under ``root``; anything that
    escapes the root raises ``ValueError`` before touching disk.
    """

    def __init__(self, root: str | Path) -> None:
        self.root = Path(root).resolve()
        self.root.mkdir(parents=True, exist_ok=True)

    # --- path safety (ported from artifact_store._validate_under) -------------------
    def _resolve(self, relative_path: str) -> Path:
        if not relative_path or not isinstance(relative_path, str):
            raise ValueError("Invalid path")
        if not _SAFE_REL.match(relative_path):
            raise ValueError("Path must be a safe relative path under the contract root")
        p = Path(relative_path)
        if p.is_absolute() or ".." in p.parts:
            raise ValueError("Path traversal not allowed")
        full = (self.root / p).resolve()
        if full != self.root:
            try:
                full.relative_to(self.root)
            except ValueError as e:
                raise ValueError("Path escapes the contract root") from e
        return full

    # --- .docx --------------------------------------------------------------------
    def read_docx(self, path: str) -> str:
        """Return the document's paragraphs (and table cells) as indexed lines.

        The index is the anchor the editing tools use, so the model reads first, then edits
        by index or by a unique snippet. Tables are appended after paragraphs.
        """
        full = self._resolve(path)
        if not full.exists():
            return f"[File not found: {path}]"
        try:
            doc = Document(str(full))
        except Exception as e:  # wrong format (e.g. a .pdf) — tell the model, don't crash the loop
            return f"[read_docx error: {path} is not a valid .docx ({e}). Use the right tool for its type.]"
        lines: list[str] = []
        for i, para in enumerate(doc.paragraphs):
            style = para.style.name if para.style else ""
            tag = f" «{style}»" if style and style != "Normal" else ""
            lines.append(f"[p{i}]{tag} {para.text}")
        for ti, table in enumerate(doc.tables):
            for ri, row in enumerate(table.rows):
                cells = " | ".join(c.text for c in row.cells)
                lines.append(f"[t{ti}r{ri}] {cells}")
        body = "\n".join(lines)
        return body or "[empty document]"

    @staticmethod
    def _set_paragraph_text(para, text: str) -> None:
        """Replace a paragraph's text while keeping the first run's formatting.

        python-docx splits a paragraph into runs (formatting spans). Rewriting the first run
        and dropping the rest preserves the leading run's font/size/bold rather than flattening
        to default — good enough for clause-level edits in the POC.
        """
        if para.runs:
            para.runs[0].text = text
            for extra in para.runs[1:]:
                extra._element.getparent().remove(extra._element)
        else:
            para.add_run(text)

    def edit_docx_paragraph(self, path: str, index: int, new_text: str) -> str:
        """Replace paragraph ``index`` (the ``[pN]`` from read_docx) with ``new_text``."""
        full = self._resolve(path)
        if not full.exists():
            return f"[File not found: {path}]"
        doc = Document(str(full))
        paras = doc.paragraphs
        if index < 0 or index >= len(paras):
            return f"[edit error: paragraph index {index} out of range (0..{len(paras) - 1})]"
        before = paras[index].text
        self._set_paragraph_text(paras[index], new_text)
        doc.save(str(full))
        return f"Edited [p{index}] of {path}.\n- before: {before!r}\n- after:  {new_text!r}"

    def replace_docx_text(self, path: str, old: str, new: str) -> str:
        """Anchored clause edit: replace a unique paragraph snippet (edit_file semantics).

        ``old`` must match the text of exactly one paragraph (substring). Fails loudly on
        zero or multiple matches so the model disambiguates instead of clobbering.
        """
        full = self._resolve(path)
        if not full.exists():
            return f"[File not found: {path}]"
        if not old:
            return "[edit error: 'old' must be non-empty]"
        doc = Document(str(full))
        hits = [p for p in doc.paragraphs if old in p.text]
        if not hits:
            return f"[edit error: snippet not found: {old!r}]"
        if len(hits) > 1:
            return (f"[edit error: snippet matched {len(hits)} paragraphs; add context to make it "
                    f"unique, or use edit_docx_paragraph with an index]")
        para = hits[0]
        para_new = para.text.replace(old, new)
        self._set_paragraph_text(para, para_new)
        doc.save(str(full))
        return f"Replaced clause text in {path}.\n- {old!r} → {new!r}"

    # --- .pptx --------------------------------------------------------------------
    def read_pptx(self, path: str) -> str:
        """Return slide/shape text as indexed lines (``[sN.M]`` = slide N, shape M)."""
        full = self._resolve(path)
        if not full.exists():
            return f"[File not found: {path}]"
        try:
            prs = Presentation(str(full))
        except Exception as e:  # wrong format (e.g. a .pdf) — tell the model, don't crash the loop
            return f"[read_pptx error: {path} is not a valid .pptx ({e}). Use the right tool for its type.]"
        lines: list[str] = []
        for si, slide in enumerate(prs.slides):
            for mi, shape in enumerate(slide.shapes):
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    lines.append(f"[s{si}.{mi}] {text}")
        body = "\n".join(lines)
        return body or "[no text shapes]"

    # --- .pdf (image-only / flattened posters) ------------------------------------
    #
    # A flattened PDF has NO text layer — the title is pixels in an embedded image. So we can't
    # edit text; we remove the old title by content-aware inpainting (reconstructs the paper
    # texture so it blends, instead of pasting a flat box) and typeset the new title in a serif
    # font. Output is a NEW file; the source is never overwritten. Fidelity is approximate.

    @staticmethod
    def _serif_font(size: int):
        """Return a serif TrueType font, discovered at runtime (no hardcoded path).

        Tries fontconfig (``fc-match``) then a few common locations; falls back to PIL's default
        bitmap font if nothing serif is installed. Keeps the tool portable across machines.
        """
        import shutil
        import subprocess

        from PIL import ImageFont

        candidates: list[str] = []
        fc = shutil.which("fc-match")
        if fc:
            try:
                out = subprocess.run([fc, "-f", "%{file}", "serif"], capture_output=True,
                                     text=True, timeout=5).stdout.strip()
                if out:
                    candidates.append(out)
            except Exception:
                pass
        candidates += [
            "/usr/share/fonts/truetype/dejavu/DejaVuSerif.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationSerif-Regular.ttf",
            "C:/Windows/Fonts/times.ttf", "C:/Windows/Fonts/georgia.ttf",
            "/System/Library/Fonts/Supplemental/Times New Roman.ttf",
        ]
        for path in candidates:
            try:
                return ImageFont.truetype(path, size)
            except Exception:
                continue
        return ImageFont.load_default()

    @staticmethod
    def _ocr_band(img, top: int, bottom: int) -> str:
        """OCR a horizontal band of the page image to read back its rendered text.

        Used to verify a raster retitle: the new title is pixels, so the only way to confirm it
        (or read the original) is to recognise it. Returns "" if OCR is unavailable.
        """
        try:
            import pytesseract
        except Exception:
            return ""
        crop = img.crop((0, max(0, top - 6), img.width, bottom + 6))
        try:
            text = pytesseract.image_to_string(crop)
        except Exception:
            return ""
        return " ".join(text.split())

    @staticmethod
    def _detect_title_band(gray) -> tuple[int, int]:
        """Find the topmost large-font dark text band (the title) by row profiling.

        Returns ``(top, bottom)`` in image pixels, with the bottom extended to swallow descender
        tails but clamped before the next (smaller) text band so the body text is left intact.
        """
        import numpy as np

        H = gray.shape[0]
        rowfrac = (gray < 110).mean(1)
        bands, on, start = [], False, 0
        for y in range(int(H * 0.5)):
            lit = rowfrac[y] > 0.04
            if lit and not on:
                start, on = y, True
            elif not lit and on:
                bands.append((start, y))
                on = False
        # Title = first band(s) tall enough to be the big heading (skip thin grunge-frame bands).
        big = [(s, e) for s, e in bands if (e - s) >= 30 and s > int(H * 0.05)]
        if not big:
            raise ValueError("could not detect a title band on page 1")
        top = big[0][0]
        last_big_bottom = big[-1][1] if (big[-1][0] - big[0][0]) < H * 0.3 else big[0][1]
        # Extend past descenders ("g","p" tails sit within ~28px of the baseline and may register
        # as their own faint band) but stop before the real body line further below.
        body_starts = [s for s, _ in bands if s > last_big_bottom + 28]
        ceiling = (body_starts[0] - 2) if body_starts else last_big_bottom + 30
        bottom = min(last_big_bottom + 30, ceiling)
        return max(0, top - 12), bottom

    def _current_pdf(self, path: str) -> str:
        """Resolve the live revision: prefer the ``*_edited.pdf`` sibling if retitle_pdf produced one.

        This is what makes the reviewer gate work without hardcoding paths — a reviewer that
        reads the original path automatically sees the latest edited version.
        """
        edited = path.rsplit(".", 1)[0] + "_edited.pdf"
        try:
            if self._resolve(edited).exists():
                return edited
        except ValueError:
            pass
        return path

    def read_pdf(self, path: str) -> str:
        """Report page 1: text-layer status, image dims, detected title band, and OCR'd title text.

        Reads the current revision (``*_edited.pdf`` if present) and OCRs the title band so the
        actual rendered title can be verified — essential for an image-only PDF where the title
        is pixels, not text.
        """
        import io

        import fitz  # PyMuPDF
        import numpy as np
        from PIL import Image

        live = self._current_pdf(path)
        full = self._resolve(live)
        if not full.exists():
            return f"[File not found: {path}]"
        note = "" if live == path else f" (reading latest revision '{live}')"
        doc = fitz.open(str(full))
        page = doc[0]
        text = page.get_text().strip()
        if text:
            return (f"{path}{note}: {len(doc)} page(s). HAS a text layer ({len(text)} chars) — prefer "
                    f"a real PDF-text editor.\nFirst chars: {text[:200]!r}")
        imgs = page.get_images(full=True)
        if not imgs:
            return f"{path}: no text layer and no image — nothing to retitle."
        base = doc.extract_image(imgs[0][0])
        img = Image.open(io.BytesIO(base["image"])).convert("RGB")
        gray = np.asarray(img).mean(2)
        try:
            t, b = self._detect_title_band(gray)
            ocr = self._ocr_band(img, t, b)
            title = f"current title text (OCR): {ocr!r}" if ocr else "OCR unavailable"
            band = f"detected title band rows {t}..{b}; {title}"
        except ValueError as e:
            band = f"title band NOT auto-detected ({e}); pass top/bottom explicitly"
        return (f"{path}{note}: {len(doc)} page(s). NO text layer — flattened image "
                f"{img.width}x{img.height}. {band}. Use retitle_pdf to replace the title.")

    def retitle_pdf(self, path: str, new_title: str, top: int | None = None,
                    bottom: int | None = None) -> str:
        """Replace the title of a flattened (image-only) PDF; write ``*(edited).pdf``.

        Inpaints the old title (blends with the paper texture) then typesets ``new_title`` in a
        serif face, wrapped to two lines. ``top``/``bottom`` override the auto-detected band.
        """
        import io

        import cv2
        import fitz
        import numpy as np
        from PIL import Image, ImageDraw, ImageFont

        full = self._resolve(path)
        if not full.exists():
            return f"[File not found: {path}]"
        doc = fitz.open(str(full))
        page = doc[0]
        imgs = page.get_images(full=True)
        if not imgs:
            return "[retitle_pdf error: page 1 has no embedded image]"
        base = doc.extract_image(imgs[0][0])
        img = Image.open(io.BytesIO(base["image"])).convert("RGB")
        arr = np.asarray(img).copy()
        gray = arr.mean(2)
        W = img.width

        if top is None or bottom is None:
            t, b = self._detect_title_band(gray)
        else:
            t, b = int(top), int(bottom)
        # Horizontal band: stay off the grunge frame (clamp to central 8%..92%).
        l, r = int(W * 0.08), int(W * 0.92)

        # Inpaint old title: mask dark glyphs (incl. faint descenders), clamp inside the band so
        # dilation can't spill onto the body text just below.
        mask = np.zeros(gray.shape, np.uint8)
        region = gray[t:b, l:r] < 135
        mask[t:b, l:r][region] = 255
        mask = cv2.dilate(mask, np.ones((5, 5), np.uint8), iterations=2)
        clamp = np.zeros_like(mask)
        clamp[t:b, l:r] = mask[t:b, l:r]
        bgr = cv2.cvtColor(arr, cv2.COLOR_RGB2BGR)
        bgr = cv2.inpaint(bgr, clamp, 3, cv2.INPAINT_NS)
        img = Image.fromarray(cv2.cvtColor(bgr, cv2.COLOR_BGR2RGB))

        # Typeset the new title, wrapped to <=2 lines, auto-fit to the band.
        draw = ImageDraw.Draw(img)
        words = new_title.split()
        mid = len(words) // 2 or 1
        lines = [new_title] if len(words) <= 3 else [" ".join(words[:mid]), " ".join(words[mid:])]
        box_w, box_h = r - l, b - t
        size = 8
        font = self._serif_font(size)
        while size < 200:
            f = self._serif_font(size + 2)
            widest = max(draw.textlength(ln, font=f) for ln in lines)
            line_h = (f.getbbox("Hg")[3] - f.getbbox("Hg")[1]) * 1.25
            if widest > box_w * 0.98 or line_h * len(lines) > box_h * 0.98:
                break
            size, font = size + 2, f
        line_h = (font.getbbox("Hg")[3] - font.getbbox("Hg")[1]) * 1.25
        y = t + (box_h - line_h * len(lines)) / 2
        for ln in lines:
            x = l + (box_w - draw.textlength(ln, font=font)) / 2
            draw.text((x, y), ln, fill=(20, 20, 20), font=font)
            y += line_h

        out_rel = path.rsplit(".", 1)[0] + "_edited.pdf"
        out_full = self._resolve(out_rel)
        out = fitz.open()
        pg = out.new_page(width=page.rect.width, height=page.rect.height)
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=92)
        pg.insert_image(pg.rect, stream=buf.getvalue())
        out.save(str(out_full))
        # OCR the result so the caller (and the reviewer gate) can verify the rendered title.
        readback = self._ocr_band(img, t, b)
        verify = f" OCR readback: {readback!r}." if readback else ""
        return (f"Retitled {path} → {out_rel}. Inpainted band rows {t}..{b}, "
                f"typeset {len(lines)} line(s) at {size}px serif.{verify}")

    def edit_pptx_text(self, path: str, slide: int, shape: int, new_text: str) -> str:
        """Replace the text of slide ``slide`` shape ``shape`` (the ``[sN.M]`` anchor)."""
        full = self._resolve(path)
        if not full.exists():
            return f"[File not found: {path}]"
        prs = Presentation(str(full))
        slides = list(prs.slides)
        if slide < 0 or slide >= len(slides):
            return f"[edit error: slide {slide} out of range (0..{len(slides) - 1})]"
        shapes = list(slides[slide].shapes)
        if shape < 0 or shape >= len(shapes):
            return f"[edit error: shape {shape} out of range on slide {slide}]"
        target = shapes[shape]
        if not target.has_text_frame:
            return f"[edit error: shape [s{slide}.{shape}] has no text frame]"
        tf = target.text_frame
        before = tf.text
        # Keep the first paragraph/run formatting; drop extras (POC-level fidelity).
        tf.paragraphs[0].runs[0].text = new_text if tf.paragraphs[0].runs else None
        if not tf.paragraphs[0].runs:
            tf.paragraphs[0].add_run().text = new_text
        for p in tf.paragraphs[1:]:
            p._p.getparent().remove(p._p)
        prs.save(str(full))
        return f"Edited [s{slide}.{shape}] of {path}.\n- before: {before!r}\n- after: {new_text!r}"
